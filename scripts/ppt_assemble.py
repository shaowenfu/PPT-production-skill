#!/usr/bin/env python
"""Assemble final PPTX from AI-generated slide images.

This is Program 7 in the PPT workflow. It reads the slide plan and assets
to assemble the final PowerPoint presentation by inserting each generated
image as a full-page slide.

Usage:
    python scripts/ppt_assemble.py --project-dir PPT/03 [--output-json]
"""

from __future__ import annotations

import argparse
import sys
from datetime import datetime
from pathlib import Path
from typing import Any

REPO_ROOT = Path(__file__).resolve().parents[1]
if str(REPO_ROOT) not in sys.path:
    sys.path.insert(0, str(REPO_ROOT))

from pptx import Presentation
from pptflow.cli import run_cli
from pptflow.errors import InputError, StateStoreError
from pptflow.json_io import read_json
from pptflow.paths import ProjectPaths
from pptflow.ppt_builder import (
    create_presentation,
    get_blank_layout,
    SLIDE_WIDTH,
    SLIDE_HEIGHT,
    add_speaker_notes
)
from pptflow.schemas import AssetManifest, SlidePlanDocument, SlideDraftDocument
from pptflow.state_store import (
    append_transition,
    load_state,
    save_state,
    set_artifact,
)

TOOL_NAME = "ppt_assemble"

def _add_script_args(parser: argparse.ArgumentParser) -> argparse.ArgumentParser:
    parser.add_argument("--project-dir", required=True, help="项目目录")
    parser.add_argument("--output-json", action="store_true", default=False, help="结果输出到 stdout")
    return parser

def handle_assemble(args: argparse.Namespace) -> dict[str, Any]:
    project_dir = Path(args.project_dir).expanduser().resolve()
    
    # 1. 加载状态与路径
    state = load_state(project_dir)
    project_id = state["project_id"]
    
    # 2. 读取规划与资产清单
    plan_path = project_dir / "plan" / "plan.json"
    manifest_path = project_dir / "assets" / "manifest.json"
    draft_path = project_dir / "draft" / "slide_draft.json"
    
    if not plan_path.exists() or not manifest_path.exists():
        raise InputError("规划文件或资产清单缺失，请先执行前面的步骤")

    plan = SlidePlanDocument(**read_json(plan_path))
    manifest = AssetManifest(**read_json(manifest_path))
    
    # 可选：读取文案用于填充演讲者备注
    draft = None
    if draft_path.exists():
        draft = SlideDraftDocument(**read_json(draft_path))
        draft_map = {s.page_id: s for s in draft.slides}
    else:
        draft_map = {}

    # 3. 构建资产映射
    asset_map = {item.page_id: project_dir / item.file_path for item in manifest.items}

    # 4. 创建 PPT 并按顺序插入图片
    prs = create_presentation()
    blank_layout = get_blank_layout(prs)
    
    images_count = 0
    for page in plan.pages:
        slide = prs.slides.add_slide(blank_layout)
        image_path = asset_map.get(page.page_id)
        
        if image_path and image_path.exists():
            # 将图片插入为全屏
            slide.shapes.add_picture(
                str(image_path),
                0, 0,
                SLIDE_WIDTH, SLIDE_HEIGHT
            )
            images_count += 1
        
        # 填充演讲者备注（辅助演示）
        if page.page_id in draft_map:
            add_speaker_notes(slide, draft_map[page.page_id].speaker_note)

    # 5. 保存结果
    deck_dir = project_dir / "deck"
    deck_dir.mkdir(exist_ok=True)
    output_pptx = deck_dir / "deck.pptx"
    prs.save(str(output_pptx))

    # 6. 更新状态
    state = set_artifact(state, "deck", "deck/deck.pptx", exists=True)
    state["current_state"] = "DeckAssembled"
    state["last_completed_step"] = TOOL_NAME
    
    append_transition(state, {
        "timestamp": datetime.now().astimezone().isoformat(timespec="seconds"),
        "from_state": state.get("current_state", "AssetsGenerated"),
        "to_state": "DeckAssembled",
        "trigger": "tool_success",
        "step": TOOL_NAME,
        "note": f"Assembled {len(plan.pages)} slides using {images_count} full-page images.",
    })
    save_state(project_dir, state)

    return {
        "project_id": project_id,
        "total_slides": len(plan.pages),
        "images_inserted": images_count,
        "output_file": "deck/deck.pptx"
    }

def main() -> int:
    parser = argparse.ArgumentParser(prog=TOOL_NAME)
    _add_script_args(parser)
    return run_cli(handle_assemble, tool=TOOL_NAME, parser=parser)

if __name__ == "__main__":
    raise SystemExit(main())
