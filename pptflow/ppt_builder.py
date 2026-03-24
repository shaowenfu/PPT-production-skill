"""PowerPoint generation module using python-pptx.

Minimalist version: Focusing on image-as-slide workflow.
"""

from __future__ import annotations

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches

# Slide dimensions for 16:9 aspect ratio (in EMUs - English Metric Units)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

def create_presentation() -> Presentation:
    """Create a new blank presentation with 16:9 aspect ratio."""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs

def get_blank_layout(prs: Presentation):
    """Get the blank slide layout from the presentation."""
    for layout in prs.slide_layouts:
        if layout.name == "Blank" or "blank" in layout.name.lower():
            return layout
    return prs.slide_layouts[-1] if len(prs.slide_layouts) > 0 else prs.slide_layouts[0]

def add_speaker_notes(slide, note: str) -> None:
    """Add speaker notes to a slide."""
    if note:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.text = note

__all__ = [
    "create_presentation",
    "get_blank_layout",
    "SLIDE_WIDTH",
    "SLIDE_HEIGHT",
    "add_speaker_notes"
]
